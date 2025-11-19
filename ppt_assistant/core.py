import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from .config import DEFAULT_CONFIG


class PPTAssistant:
    """PPT 助手類"""

    def __init__(self, json_file, ppt_config=None):
        """
        初始化 PPT 助手

        Args:
            json_file: JSON 配置文件路徑
            ppt_config: PPT 配置對象，默認為 DEFAULT_CONFIG
        """
        with open(json_file, 'r', encoding='utf-8') as f:
            self.json_config = json.load(f)

        self.ppt_config = ppt_config or DEFAULT_CONFIG

        self.template_path = self.json_config.get('template')
        self.output_path = self.json_config.get('output')
        self.slides_data = self.json_config.get('slides', [])

        # 載入模板或創建新簡報
        if self.template_path and os.path.exists(self.template_path):
            self.prs = Presentation(self.template_path)
        else:
            self.prs = Presentation()

        # 設置為 16:9 比例 (PowerPoint 標準寬屏尺寸)
        self.prs.slide_width = Inches(self.ppt_config.SLIDE_WIDTH)
        self.prs.slide_height = Inches(self.ppt_config.SLIDE_HEIGHT)

    def create_presentation(self):
        """創建完整的簡報"""
        for slide_data in self.slides_data:
            layout_type = slide_data.get('layout')
            content = slide_data.get('content', {})

            if layout_type == 'title':
                self._add_title_slide(content)
            elif layout_type == 'content':
                self._add_content_slide(content)
            elif layout_type == 'two_column':
                self._add_two_column_slide(content)
            elif layout_type == 'image':
                self._add_image_slide(content)
            elif layout_type == 'table':
                self._add_table_slide(content)
            elif layout_type == 'chart':
                self._add_chart_slide(content)
            else:
                print(f"警告: 不支援的佈局類型 '{layout_type}'")

        # 保存簡報
        self.prs.save(self.output_path)
        print(f"簡報已成功生成: {self.output_path}")

    def _add_title_slide(self, content):
        """添加標題頁"""
        slide_layout = self.prs.slide_layouts[self.ppt_config.TITLE_LAYOUT_INDEX]
        slide = self.prs.slides.add_slide(slide_layout)

        # 設置標題
        if slide.shapes.title:
            title = slide.shapes.title
            title.text = content.get('title', '')
            title.left = Inches(self.ppt_config.MARGIN)
            title.top = Inches(self.ppt_config.TITLE_TOP_POSITION)
            title.width = Inches(self.ppt_config.content_width)
            title.height = Inches(self.ppt_config.TITLE_HEIGHT_SIZE)
            tf = title.text_frame
            tf.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            tf.paragraphs[0].font.size = Pt(self.ppt_config.TITLE_FONT_SIZE)
            tf.paragraphs[0].font.bold = True

        # 設置副標題
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = content.get('subtitle', '')
            subtitle.left = Inches(self.ppt_config.MARGIN)
            subtitle.top = Inches(self.ppt_config.SUBTITLE_TOP_POSITION)
            subtitle.width = Inches(self.ppt_config.content_width)
            subtitle.height = Inches(self.ppt_config.SUBTITLE_HEIGHT_SIZE)
            tf = subtitle.text_frame
            tf.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            tf.paragraphs[0].font.size = Pt(self.ppt_config.SUBTITLE_FONT_SIZE)

    def _add_content_slide(self, content):
        """添加內容頁"""
        layout_index = self.ppt_config.CONTENT_LAYOUT_INDEX if len(self.prs.slide_layouts) > self.ppt_config.CONTENT_LAYOUT_INDEX else self.ppt_config.CONTENT_LAYOUT_FALLBACK_INDEX
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加標題
        title_shape = slide.shapes.add_textbox(
            Inches(self.ppt_config.MARGIN),
            Inches(self.ppt_config.SLIDE_TITLE_TOP),
            Inches(self.ppt_config.content_width),
            Inches(self.ppt_config.SLIDE_TITLE_HEIGHT)
        )
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = content.get('title', '')
        p.font.size = Pt(self.ppt_config.SLIDE_TITLE_FONT_SIZE)
        p.font.bold = True
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

        # 添加內容
        body_shape = slide.shapes.add_textbox(
            Inches(self.ppt_config.MARGIN),
            Inches(self.ppt_config.content_top),
            Inches(self.ppt_config.content_width),
            Inches(self.ppt_config.content_height)
        )
        tf = body_shape.text_frame
        tf.word_wrap = True

        body_text = content.get('body', '')
        if isinstance(body_text, list):
            for idx, item in enumerate(body_text):
                if idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = str(item)
                p.level = 0
                p.font.size = Pt(self.ppt_config.BODY_FONT_SIZE)
                p.space_before = Pt(self.ppt_config.PARAGRAPH_SPACE_BEFORE)
                p.space_after = Pt(self.ppt_config.PARAGRAPH_SPACE_AFTER)
                p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        else:
            p = tf.paragraphs[0]
            p.text = str(body_text)
            p.font.size = Pt(self.ppt_config.BODY_FONT_SIZE)
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    def _add_two_column_slide(self, content):
        """添加雙欄佈局頁"""
        layout_index = self.ppt_config.CONTENT_LAYOUT_INDEX if len(self.prs.slide_layouts) > self.ppt_config.CONTENT_LAYOUT_INDEX else self.ppt_config.CONTENT_LAYOUT_FALLBACK_INDEX
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加標題
        title_shape = slide.shapes.add_textbox(
            Inches(self.ppt_config.MARGIN),
            Inches(self.ppt_config.SLIDE_TITLE_TOP),
            Inches(self.ppt_config.content_width),
            Inches(self.ppt_config.SLIDE_TITLE_HEIGHT)
        )
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = content.get('title', '')
        p.font.size = Pt(self.ppt_config.SLIDE_TITLE_FONT_SIZE)
        p.font.bold = True
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

        # 計算雙欄尺寸
        column_width = (self.ppt_config.content_width - self.ppt_config.COLUMN_GAP) / 2

        # 左欄
        left_box = slide.shapes.add_textbox(
            Inches(self.ppt_config.MARGIN),
            Inches(self.ppt_config.content_top),
            Inches(column_width),
            Inches(self.ppt_config.content_height)
        )
        left_tf = left_box.text_frame
        left_tf.word_wrap = True
        left_p = left_tf.paragraphs[0]
        left_p.text = str(content.get('left', ''))
        left_p.font.size = Pt(self.ppt_config.COLUMN_FONT_SIZE)
        left_p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        left_p.space_after = Pt(self.ppt_config.COLUMN_SPACE_AFTER)

        # 右欄
        right_left_pos = self.ppt_config.MARGIN + column_width + self.ppt_config.COLUMN_GAP
        right_box = slide.shapes.add_textbox(
            Inches(right_left_pos),
            Inches(self.ppt_config.content_top),
            Inches(column_width),
            Inches(self.ppt_config.content_height)
        )
        right_tf = right_box.text_frame
        right_tf.word_wrap = True
        right_p = right_tf.paragraphs[0]
        right_p.text = str(content.get('right', ''))
        right_p.font.size = Pt(self.ppt_config.COLUMN_FONT_SIZE)
        right_p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        right_p.space_after = Pt(self.ppt_config.COLUMN_SPACE_AFTER)

    def _add_image_slide(self, content):
        """添加圖片頁"""
        layout_index = self.ppt_config.CONTENT_LAYOUT_INDEX if len(self.prs.slide_layouts) > self.ppt_config.CONTENT_LAYOUT_INDEX else self.ppt_config.CONTENT_LAYOUT_FALLBACK_INDEX
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加標題
        title_shape = slide.shapes.add_textbox(
            Inches(self.ppt_config.MARGIN),
            Inches(self.ppt_config.SLIDE_TITLE_TOP),
            Inches(self.ppt_config.content_width),
            Inches(self.ppt_config.SLIDE_TITLE_HEIGHT)
        )
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = content.get('title', '')
        p.font.size = Pt(self.ppt_config.SLIDE_TITLE_FONT_SIZE)
        p.font.bold = True
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

        # 添加圖片
        image_path = content.get('image', '')
        if image_path and os.path.exists(image_path):
            # 圖片居中,使用大部分內容區域
            img_width = self.ppt_config.content_width * self.ppt_config.IMAGE_WIDTH_RATIO
            img_left = self.ppt_config.MARGIN + (self.ppt_config.content_width - img_width) / 2
            img_top = self.ppt_config.content_top + self.ppt_config.IMAGE_TOP_OFFSET
            slide.shapes.add_picture(
                image_path,
                Inches(img_left),
                Inches(img_top),
                width=Inches(img_width)
            )

        # 添加圖片說明(可選)
        caption = content.get('caption', '')
        if caption:
            caption_shape = slide.shapes.add_textbox(
                Inches(self.ppt_config.MARGIN),
                Inches(self.ppt_config.SLIDE_HEIGHT - self.ppt_config.CAPTION_TOP_FROM_BOTTOM),
                Inches(self.ppt_config.content_width),
                Inches(self.ppt_config.CAPTION_HEIGHT)
            )
            caption_tf = caption_shape.text_frame
            caption_p = caption_tf.paragraphs[0]
            caption_p.text = caption
            caption_p.font.size = Pt(self.ppt_config.CAPTION_FONT_SIZE)
            caption_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            caption_p.font.italic = True

    def _add_table_slide(self, content):
        """添加表格頁"""
        layout_index = self.ppt_config.CONTENT_LAYOUT_INDEX if len(self.prs.slide_layouts) > self.ppt_config.CONTENT_LAYOUT_INDEX else self.ppt_config.CONTENT_LAYOUT_FALLBACK_INDEX
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加標題
        title_shape = slide.shapes.add_textbox(
            Inches(self.ppt_config.MARGIN),
            Inches(self.ppt_config.SLIDE_TITLE_TOP),
            Inches(self.ppt_config.content_width),
            Inches(self.ppt_config.SLIDE_TITLE_HEIGHT)
        )
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = content.get('title', '')
        p.font.size = Pt(self.ppt_config.SLIDE_TITLE_FONT_SIZE)
        p.font.bold = True
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

        # 獲取表格數據
        table_data = content.get('table', {})
        headers = table_data.get('headers', [])
        rows = table_data.get('rows', [])

        if headers and rows:
            # 計算表格尺寸
            rows_count = len(rows) + 1  # +1 for header
            cols_count = len(headers)

            # 添加表格 - 使用完整內容區域
            table_width = self.ppt_config.content_width * self.ppt_config.TABLE_WIDTH_RATIO
            table_left = self.ppt_config.MARGIN + (self.ppt_config.content_width - table_width) / 2

            table_shape = slide.shapes.add_table(
                rows_count,
                cols_count,
                Inches(table_left),
                Inches(self.ppt_config.content_top),
                Inches(table_width),
                Inches(self.ppt_config.content_height * self.ppt_config.TABLE_HEIGHT_RATIO)
            )
            table = table_shape.table

            # 填充表頭
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = str(header)
                # 格式化表頭
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(self.ppt_config.TABLE_HEADER_FONT_SIZE)
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                # 設置表頭背景色
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*self.ppt_config.TABLE_HEADER_BG_COLOR)
                # 設置文字顏色為白色
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = RGBColor(*self.ppt_config.TABLE_HEADER_TEXT_COLOR)

            # 填充數據行
            for row_idx, row_data in enumerate(rows):
                for col_idx, cell_data in enumerate(row_data):
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(cell_data)
                    # 格式化數據單元格
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(self.ppt_config.TABLE_DATA_FONT_SIZE)
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    # 交替行背景色
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*self.ppt_config.TABLE_ALT_ROW_BG_COLOR)

    def _add_chart_slide(self, content):
        """添加圖表頁"""
        layout_index = self.ppt_config.CONTENT_LAYOUT_INDEX if len(self.prs.slide_layouts) > self.ppt_config.CONTENT_LAYOUT_INDEX else self.ppt_config.CONTENT_LAYOUT_FALLBACK_INDEX
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加標題
        title_shape = slide.shapes.add_textbox(
            Inches(self.ppt_config.MARGIN),
            Inches(self.ppt_config.SLIDE_TITLE_TOP),
            Inches(self.ppt_config.content_width),
            Inches(self.ppt_config.SLIDE_TITLE_HEIGHT)
        )
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = content.get('title', '')
        p.font.size = Pt(self.ppt_config.SLIDE_TITLE_FONT_SIZE)
        p.font.bold = True
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

        # 獲取圖表數據
        chart_config = content.get('chart', {})
        chart_type = chart_config.get('type', 'bar')
        chart_data_config = chart_config.get('data', {})

        categories = chart_data_config.get('categories', [])
        series_list = chart_data_config.get('series', [])

        if categories and series_list:
            # 創建圖表數據
            chart_data = CategoryChartData()
            chart_data.categories = categories

            for series in series_list:
                series_name = series.get('name', '')
                series_values = series.get('values', [])
                chart_data.add_series(series_name, series_values)

            # 確定圖表類型
            chart_type_attr = getattr(XL_CHART_TYPE, self.ppt_config.CHART_TYPE_MAPPING.get(chart_type, 'COLUMN_CLUSTERED'))

            # 添加圖表 - 使用完整內容區域
            chart_width = self.ppt_config.content_width * self.ppt_config.CHART_WIDTH_RATIO
            chart_left = self.ppt_config.MARGIN + (self.ppt_config.content_width - chart_width) / 2

            chart_shape = slide.shapes.add_chart(
                chart_type_attr,
                Inches(chart_left),
                Inches(self.ppt_config.content_top),
                Inches(chart_width),
                Inches(self.ppt_config.content_height * self.ppt_config.CHART_HEIGHT_RATIO),
                chart_data
            )
            chart = chart_shape.chart

            # 設置圖表樣式
            chart.has_legend = True
            chart.legend.position = self.ppt_config.CHART_LEGEND_POSITION
            chart.legend.include_in_layout = False